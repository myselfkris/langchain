"""
document_converter.py
----------------------
A modular and optimized document-to-text converter.
Supports: .docx, .pdf, .pptx, .txt
"""

import os
from docx import Document
from pptx import Presentation
import PyPDF2

# ===============================
# 🧾 1. DOCX Conversion
# ===============================
def convert_docx_to_text(docx_path: str) -> str:
    """
    Extracts text from a .docx (Word) file efficiently.
    - Removes empty lines and trims spaces.
    """
    doc = Document(docx_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


# ===============================
# 📄 2. PDF Conversion (PyPDF2)
# ===============================
def convert_pdf_to_text(pdf_path: str) -> str:
    """
    Extracts text from a PDF file using PyPDF2.
    - Works best for digitally-created PDFs (not scanned ones).
    - Uses streaming (no heavy memory load).
    """
    text = []
    with open(pdf_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page_num, page in enumerate(reader.pages):
            content = page.extract_text()
            if content:
                text.append(content.strip())
    return "\n".join(text)


# ===============================
# 🖼️ 3. PPTX Conversion
# ===============================
def convert_pptx_to_text(pptx_path: str) -> str:
    """
    Extracts text from PowerPoint (.pptx) slides.
    - Reads text from all text-containing shapes.
    - Skips empty shapes for cleaner output.
    """
    prs = Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())
    return "\n".join(text)


# ===============================
# 🧾 4. TXT Conversion
# ===============================
def convert_txt_to_text(txt_path: str) -> str:
    """
    Reads and cleans plain text files.
    - Strips extra spaces and blank lines.
    """
    with open(txt_path, 'r', encoding='utf-8') as file:
        lines = [line.strip() for line in file if line.strip()]
    return "\n".join(lines)


# ===============================
# ⚙️ 5. Auto-Detect & Route
# ===============================
def convert_to_text(file_path: str) -> str:
    """
    Detects the file type based on extension
    and routes it to the correct conversion function.
    """
    ext = os.path.splitext(file_path)[1].lower()
    converters = {
        ".docx": convert_docx_to_text,
        ".pdf": convert_pdf_to_text,
        ".pptx": convert_pptx_to_text,
        ".txt": convert_txt_to_text,
    }

    if ext not in converters:
        raise ValueError(f"Unsupported file format: {ext}")

    print(f"🔍 Extracting text from: {file_path}")
    text = converters[ext](file_path)
    print(f"✅ Extracted {len(text)} characters from {file_path}")
    return text


# ===============================
# 🧪 Example usage (for testing)
# ===============================
if __name__ == "__main__":
    path = input("Enter file path: ").strip()
    text = convert_to_text(path)
    print("\nPreview of extracted text:\n" + "="*40)
    print(text[:500], "...")







